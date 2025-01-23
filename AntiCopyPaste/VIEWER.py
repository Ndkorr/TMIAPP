import base64
import json
import os



def view_custom_file(custom_file_path):
    try:
        # Step 1: Read the custom file
        with open(custom_file_path, "r") as custom_file:
            data = json.load(custom_file)

        # Step 2: Extract metadata and content
        original_type = data["metadata"]["original_type"]
        file_name = data["metadata"]["file_name"]
        decoded_content = base64.b64decode(data["content"])

        # Step 3: Save the original file
        output_file_name = os.path.abspath(f"extracted_{file_name}")
        print(f"Saving file to: {output_file_name}")

        # Write binary content
        with open(output_file_name, "wb") as original_file:
            original_file.write(decoded_content)

        print(f"Original file extracted and saved as: {output_file_name}")

        # Step 4: Optional - Display the file based on type
        if original_type == "pdf":
            display_pdf(output_file_name)
        elif original_type == "xlsx":
            display_excel(output_file_name)
        elif original_type == "pptx":
            display_ppt(output_file_name)
        else:
            print(f"File type '{original_type}' is not supported for viewing.")

    except Exception as e:
        print(f"An error occurred: {e}")

# Helper function to display PDFs
def display_pdf(file_path):
    import fitz  # PyMuPDF
    pdf_document = fitz.open(file_path)
    print(f"PDF has {pdf_document.page_count} pages. Displaying first page:")
    first_page = pdf_document[0]
    print(first_page.get_text())
    pdf_document.close()

# Helper function to display Excel files
def display_excel(file_path):
    import pandas as pd
    sheet = pd.read_excel(file_path)
    print("First 5 rows of the Excel file:")
    print(sheet.head())

# Helper function to display PowerPoint files
def display_ppt(file_path):
    from pptx import Presentation
    presentation = Presentation(file_path)
    print(f"PowerPoint has {len(presentation.slides)} slides. Titles:")
    for i, slide in enumerate(presentation.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"Slide {i + 1}: {title}")

# Example usage
custom_file = r"C:\Users\HR-IT-MATTHEW-PC\try.qfs"  # Replace with your custom file path
view_custom_file(custom_file)
