import base64
import json
import os
from tkinter import Tk, filedialog, Label, Button, Text, Scrollbar, Frame, messagebox, Toplevel
import fitz  # PyMuPDF for PDFs
import pandas as pd
from pptx import Presentation

# Main GUI Viewer Class
class CustomFileViewer:
    def __init__(self, root=None):
        if root:
            self.root = root
            self.root.withdraw()
        else:
            self.root = Tk()
            self.root.withdraw()  # Hide the main Tkinter window

        # Directly open file dialog and process the file
        self.open_file()

    def open_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a Custom File",
            filetypes=(("Custom Files", "*.myext"), ("All Files", "*.*"))
        )
        if file_path:
            self.open_viewer_window(file_path)

    def open_viewer_window(self, custom_file_path):
        viewer_window = Toplevel(self.root)
        viewer_window.title("File Viewer")
        viewer_window.geometry("800x600")

        text_display = Text(viewer_window, wrap="word")
        scrollbar = Scrollbar(viewer_window, command=text_display.yview)
        text_display.configure(yscrollcommand=scrollbar.set)

        text_display.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")

        try:
             # Handle window close event
            def on_closing():
                self.root.quit()
                self.root.destroy()

            # Read and decode the .myext file
            with open(custom_file_path, "r") as custom_file:
                data = json.load(custom_file)

            # Extract metadata and content
            original_type = data["metadata"]["original_type"]
            file_name = data["metadata"]["file_name"]
            decoded_content = base64.b64decode(data["content"])

            # Save the extracted file locally
            output_file = f"extracted_{file_name}"
            with open(output_file, "wb") as original_file:
                original_file.write(decoded_content)

            # Display content based on file type
            if original_type == "pdf":
                self.display_pdf(output_file, text_display)
            elif original_type == "xlsx":
                self.display_excel(output_file, text_display)
            elif original_type == "pptx":
                self.display_ppt(output_file, text_display)
            elif original_type == "docx":
                self.display_docx(output_file, text_display)
            else:
                messagebox.showinfo("Unsupported File", f"File type '{original_type}' is not supported for viewing.")
                return
                

            viewer_window.protocol("WM_DELETE_WINDOW", on_closing)
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")
            self.root.quit()
            self.root.destroy()

    def display_pdf(self, pdf_path, text_display):
        """Display text content of a PDF file."""
        try:
            pdf_document = fitz.open(pdf_path)
            text = ""
            for page in pdf_document:
                text += page.get_text()
            text_display.delete(1.0, "end")
            text_display.insert("end", text)
            pdf_document.close()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display PDF: {e}")

    def display_excel(self, excel_path, text_display):
        """Display the first few rows of an Excel file."""
        try:
            df = pd.read_excel(excel_path)
            text_display.delete(1.0, "end")
            text_display.insert("end", df.head().to_string())
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display Excel: {e}")

    def display_ppt(self, ppt_path, text_display):
        """Display slide titles from a PowerPoint file."""
        try:
            presentation = Presentation(ppt_path)
            text_display.delete(1.0, "end")
            for i, slide in enumerate(presentation.slides):
                title = slide.shapes.title.text if slide.shapes.title else "No Title"
                text_display.insert("end", f"Slide {i + 1}: {title}\n")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display PowerPoint: {e}")

    def display_docx(self, docx_path, text_display):
        """Display content of a Word (.docx) file."""
        try:
            from docx import Document
            document = Document(docx_path)

            # Clear previous content in the text display
            text_display.delete(1.0, "end")

            # Extract and display paragraphs
            for paragraph in document.paragraphs:
                text_display.insert("end", paragraph.text + "\n")

        except Exception as e:
            messagebox.showerror("Error", f"Failed to display Word document: {e}")

# Run the GUI application
if __name__ == "__main__":
    root = Tk()
    app = CustomFileViewer(root)
    root.mainloop()