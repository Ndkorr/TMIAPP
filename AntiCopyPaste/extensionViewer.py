import base64
import json
import os
from tkinter import Tk, filedialog, Label, Button, Text, Scrollbar, Frame, messagebox
import fitz  # PyMuPDF for PDFs
import pandas as pd
from pptx import Presentation

# Main GUI Viewer Class
class CustomFileViewer:
    def __init__(self, root):
        self.root = root
        self.root.title("Custom File Viewer")
        self.root.geometry("800x600")

        # GUI Layout
        Label(root, text="Custom File Viewer", font=("Arial", 20)).pack(pady=10)
        Button(root, text="Open .myext File", command=self.open_file).pack(pady=10)

        # File info display
        self.file_info = Label(root, text="", font=("Arial", 12), fg="blue")
        self.file_info.pack()

        # Frame for content display
        self.content_frame = Frame(root)
        self.content_frame.pack(fill="both", expand=True)

        self.text_display = Text(self.content_frame, wrap="word")
        self.scrollbar = Scrollbar(self.content_frame, command=self.text_display.yview)
        self.text_display.configure(yscrollcommand=self.scrollbar.set)

        self.text_display.pack(side="left", fill="both", expand=True)
        self.scrollbar.pack(side="right", fill="y")

    def open_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a Custom File",
            filetypes=(("Custom Files", "*.myext"), ("All Files", "*.*"))
        )
        if file_path:
            self.process_file(file_path)

    def process_file(self, custom_file_path):
        try:
            # Read and decode the .myext file
            with open(custom_file_path, "r") as custom_file:
                data = json.load(custom_file)

            # Extract metadata and content
            original_type = data["metadata"]["original_type"]
            file_name = data["metadata"]["file_name"]
            decoded_content = base64.b64decode(data["content"])

            # Save the extracted file locally
            output_file = f"extracted_{file_name}"
            with open(output_file, "wb") as original_file:
                original_file.write(decoded_content)

            # Update GUI with file info
            self.file_info.config(text=f"Viewing: {output_file}")

            # Display content based on file type
            if original_type == "pdf":
                self.display_pdf(output_file)
            elif original_type == "xlsx":
                self.display_excel(output_file)
            elif original_type == "pptx":
                self.display_ppt(output_file)
            else:
                messagebox.showinfo("Unsupported File", f"File type '{original_type}' is not supported for viewing.")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")

    def display_pdf(self, pdf_path):
        """Display text content of a PDF file."""
        try:
            pdf_document = fitz.open(pdf_path)
            text = ""
            for page in pdf_document:
                text += page.get_text()
            self.text_display.delete(1.0, "end")
            self.text_display.insert("end", text)
            pdf_document.close()
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display PDF: {e}")

    def display_excel(self, excel_path):
        """Display the first few rows of an Excel file."""
        try:
            df = pd.read_excel(excel_path)
            self.text_display.delete(1.0, "end")
            self.text_display.insert("end", df.head().to_string())
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display Excel: {e}")

    def display_ppt(self, ppt_path):
        """Display slide titles from a PowerPoint file."""
        try:
            presentation = Presentation(ppt_path)
            self.text_display.delete(1.0, "end")
            for i, slide in enumerate(presentation.slides):
                title = slide.shapes.title.text if slide.shapes.title else "No Title"
                self.text_display.insert("end", f"Slide {i + 1}: {title}\n")
        except Exception as e:
            messagebox.showerror("Error", f"Failed to display PowerPoint: {e}")

# Run the GUI application
if __name__ == "__main__":
    root = Tk()
    app = CustomFileViewer(root)
    root.mainloop()
